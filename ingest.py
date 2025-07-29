import os
import shutil
from langchain_community.vectorstores import Chroma
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from pptx import Presentation

UPLOAD_DIR = os.path.abspath("uploaded_file_vectordb")
UPLOAD_CONTEXT_PATH = "uploaded_file_context.txt"

def extract_text_from_pptx(file_path):
    """Extract text content from a .pptx file."""
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
    return text.strip()

def extract_text_from_txt(file_path):
    """Extract raw text content from a .txt file."""
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read().strip()

def ingest_file(file_path):
    print("📥 Starting ingestion for file:", file_path)

    # Clean previous vectorstore
    if os.path.exists(UPLOAD_DIR):
        print("🧹 Removing old vector DB at:", UPLOAD_DIR)
        shutil.rmtree(UPLOAD_DIR, ignore_errors=True)

    # Determine file type and extract text
    if file_path.lower().endswith(".pptx"):
        text = extract_text_from_pptx(file_path)
    elif file_path.lower().endswith(".txt"):
        text = extract_text_from_txt(file_path)
    else:
        raise ValueError("Unsupported file type. Only .pptx and .txt are allowed.")

    # Split text into chunks
    splitter = CharacterTextSplitter(separator="\n", chunk_size=1000, chunk_overlap=100)
    docs = splitter.create_documents([text])

    # Generate embeddings and build Chroma vector store
    embedding = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    vectordb = Chroma.from_documents(docs, embedding_function=embedding, persist_directory=UPLOAD_DIR)
    vectordb.persist()

    # Save context path
    with open(UPLOAD_CONTEXT_PATH, "w") as f:
        f.write(file_path)

    print("✅ Ingestion complete for:", file_path)

def clear_uploaded_context():
    """Delete the vector store and context reference file."""
    if os.path.exists(UPLOAD_DIR):
        shutil.rmtree(UPLOAD_DIR, ignore_errors=True)
        print("🧹 Removed vector DB at", UPLOAD_DIR)

    if os.path.exists(UPLOAD_CONTEXT_PATH):
        os.remove(UPLOAD_CONTEXT_PATH)
        print("🧹 Removed uploaded file context path")
